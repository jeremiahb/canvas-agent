"""
agent/intake_pipeline.py

Intelligent intake pipeline for the Canvas student agent.

Every item the crawler encounters passes through seven layers before
being stored in ChromaDB:

  1. ItemClassifier  — determines source_type, content_type, intent
  2. ItemRouter      — picks the right extractor for the content type
  3. Extractor       — pulls structured text out of raw bytes / HTML
  4. Normalizer      — assembles a canonical IntakeRecord
  5. Deduplicator    — tags exact and near duplicates; avoids re-storage
  6. QualityScorer   — assigns quality (0-1) and trust (0-1) scores
  7. DecisionLogger  — records every classification decision for debugging

New content types are added by:
  a. Adding a value to ContentType / SourceType / Intent
  b. Writing a new Extractor subclass with can_handle() + extract()
  c. Registering it in ItemRouter._EXTRACTORS (before RawTextExtractor)

Nothing is ever silently discarded — unknown items fall through to
RawTextExtractor and are flagged for review.
"""

from __future__ import annotations

import csv
import hashlib
import io
import json
import logging
import mimetypes
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import datetime, timezone
from enum import Enum
from typing import Optional
from urllib.parse import parse_qs, urlencode, urlparse, urlunparse

logger = logging.getLogger(__name__)

# ------------------------------------------------------------------ #
#  Enums                                                               #
# ------------------------------------------------------------------ #


class SourceType(str, Enum):
    WEBSITE_PAGE = "website_page"
    FILE_DOWNLOAD = "file_download"
    CLOUD_DOCUMENT = "cloud_document"
    EMBEDDED_MEDIA = "embedded_media"
    ARCHIVE = "archive"
    IMAGE = "image"
    RAW_TEXT = "raw_text"
    STRUCTURED_DATASET = "structured_dataset"
    UNKNOWN = "unknown"


class ContentType(str, Enum):
    HTML = "html"
    PDF = "pdf"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    CSV = "csv"
    JSON_DATA = "json"
    IMAGE = "image"
    AUDIO = "audio"
    VIDEO = "video"
    TEXT = "text"
    UNKNOWN = "unknown"


class Intent(str, Enum):
    REFERENCE_MATERIAL = "reference_material"
    STRUCTURED_DATA = "structured_data"
    NARRATIVE_DOCUMENT = "narrative_document"
    TABULAR = "tabular"
    MEDIA_ASSET = "media_asset"
    NAVIGATION_PAGE = "navigation_page"
    INDEX_PAGE = "index_page"
    LOGIN_BARRIER = "login_barrier"
    DUPLICATE = "duplicate"
    UNKNOWN = "unknown"


# ------------------------------------------------------------------ #
#  Dataclasses                                                         #
# ------------------------------------------------------------------ #


@dataclass
class ItemContext:
    """Rich context captured at crawl time — available to all pipeline layers."""
    source_url: str
    referring_page: str = ""
    anchor_text: str = ""
    page_title: str = ""
    document_title: str = ""
    course_name: str = ""
    course_id: str = ""
    module_name: str = ""
    parent_section: str = ""
    # "module" | "syllabus" | "announcements" | "files_page" | "pages"
    source_label: str = ""
    language: str = "en"
    author: str = ""
    publisher: str = ""
    timestamp: str = ""           # ISO-8601 or empty
    # "public" | "blocked" | "login_required" | "redirect" | "unknown"
    access_status: str = "public"


@dataclass
class Classification:
    source_type: SourceType
    content_type: ContentType
    intent: Intent
    mime_type: str = ""
    confidence: float = 1.0
    reasons: list[str] = field(default_factory=list)


@dataclass
class ExtractedContent:
    raw_text: str
    cleaned_text: str
    structured: dict = field(default_factory=dict)
    # structured keys: headings[], tables[], lists[], code_blocks[],
    #                  links[], captions[], slides[], schema{}, pages[]
    page_count: int = 0
    slide_count: int = 0
    sheet_names: list[str] = field(default_factory=list)
    word_count: int = 0
    is_image_heavy: bool = False
    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


@dataclass
class IntakeRecord:
    """Fully normalized output record — one per crawled item."""
    unique_id: str
    canonical_source: str
    input_type: str           # SourceType value
    content_type: str         # ContentType value
    processing_route: str     # extractor class name used
    intent: str               # Intent value
    raw_text: str
    cleaned_text: str
    structured_elements: dict
    metadata: dict            # ItemContext fields + derived fields
    quality_score: float
    trust_score: float
    dedup_result: str         # "new" | "exact_duplicate" | "near_duplicate"
    dedup_canonical_id: str
    errors: list[str]
    warnings: list[str]
    # "search" | "analytics" | "archival" | "ignore"
    recommended_use: str
    decisions: list[dict]     # [{step, decision, reason, fallback_used, timestamp}]


# ------------------------------------------------------------------ #
#  Layer 1: ItemClassifier                                             #
# ------------------------------------------------------------------ #

# MIME → ContentType
_MIME_TO_CT: dict[str, ContentType] = {
    "text/html":                   ContentType.HTML,
    "application/xhtml+xml":       ContentType.HTML,
    "application/pdf":             ContentType.PDF,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ContentType.DOCX,
    "application/msword":          ContentType.DOCX,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ContentType.XLSX,
    "application/vnd.ms-excel":    ContentType.XLSX,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ContentType.PPTX,
    "application/vnd.ms-powerpoint": ContentType.PPTX,
    "text/csv":                    ContentType.CSV,
    "application/json":            ContentType.JSON_DATA,
    "text/plain":                  ContentType.TEXT,
}

# URL file extension → ContentType
_EXT_TO_CT: dict[str, ContentType] = {
    ".pdf":  ContentType.PDF,
    ".docx": ContentType.DOCX,
    ".doc":  ContentType.DOCX,
    ".xlsx": ContentType.XLSX,
    ".xls":  ContentType.XLSX,
    ".pptx": ContentType.PPTX,
    ".ppt":  ContentType.PPTX,
    ".csv":  ContentType.CSV,
    ".json": ContentType.JSON_DATA,
    ".txt":  ContentType.TEXT,
    ".md":   ContentType.TEXT,
    ".png":  ContentType.IMAGE,
    ".jpg":  ContentType.IMAGE,
    ".jpeg": ContentType.IMAGE,
    ".gif":  ContentType.IMAGE,
    ".webp": ContentType.IMAGE,
    ".svg":  ContentType.IMAGE,
    ".mp4":  ContentType.VIDEO,
    ".webm": ContentType.VIDEO,
    ".mov":  ContentType.VIDEO,
    ".avi":  ContentType.VIDEO,
    ".mp3":  ContentType.AUDIO,
    ".wav":  ContentType.AUDIO,
    ".ogg":  ContentType.AUDIO,
    ".m4a":  ContentType.AUDIO,
    ".zip":  ContentType.UNKNOWN,   # ArchiveExtractor handles
    ".tar":  ContentType.UNKNOWN,
    ".gz":   ContentType.UNKNOWN,
}

# Canvas URL path fragment → SourceType hint
_CANVAS_FILE_PATTERNS = ("/files/", "/download", "/file_download")
_CLOUD_DOC_HOSTS = {"docs.google.com", "drive.google.com", "onedrive.live.com",
                    "sharepoint.com", "1drv.ms"}
_MEDIA_PLATFORMS  = {"youtube.com", "youtu.be", "vimeo.com", "loom.com"}


class ItemClassifier:
    """
    Rules-based classifier. No AI calls — fast and deterministic.
    Priority: MIME type > URL extension > content sniffing > fallback.
    """

    def classify(
        self,
        url: str,
        mime_type: str = "",
        content_snippet: str = "",
        context: Optional[ItemContext] = None,
    ) -> Classification:
        reasons: list[str] = []
        content_type = ContentType.UNKNOWN
        source_type = SourceType.UNKNOWN
        intent = Intent.UNKNOWN
        confidence = 1.0

        parsed = urlparse(url)
        host = (parsed.hostname or "").lower()
        path = parsed.path.lower()
        ext = ""
        if "." in path.split("/")[-1]:
            ext = "." + path.split("/")[-1].rsplit(".", 1)[-1]

        # ── 1. MIME type (most reliable) ──────────────────────────────
        normalized_mime = mime_type.split(";")[0].strip().lower()
        if normalized_mime in _MIME_TO_CT:
            content_type = _MIME_TO_CT[normalized_mime]
            reasons.append(f"mime:{normalized_mime}")
        elif normalized_mime.startswith("image/"):
            content_type = ContentType.IMAGE
            reasons.append("mime:image/*")
        elif normalized_mime.startswith("audio/"):
            content_type = ContentType.AUDIO
            reasons.append("mime:audio/*")
        elif normalized_mime.startswith("video/"):
            content_type = ContentType.VIDEO
            reasons.append("mime:video/*")

        # ── 2. URL extension fallback ─────────────────────────────────
        if content_type == ContentType.UNKNOWN and ext in _EXT_TO_CT:
            content_type = _EXT_TO_CT[ext]
            reasons.append(f"ext:{ext}")
            if content_type == ContentType.UNKNOWN:
                content_type = ContentType.UNKNOWN   # archive — router handles

        # ── 3. Content sniffing ───────────────────────────────────────
        if content_type == ContentType.UNKNOWN and content_snippet:
            snip = content_snippet[:512]
            if snip.lstrip().startswith(("<!DOCTYPE", "<html", "<HTML")):
                content_type = ContentType.HTML
                reasons.append("sniff:html")
            elif snip.startswith("%PDF"):
                content_type = ContentType.PDF
                reasons.append("sniff:pdf-magic")
            elif snip.lstrip().startswith(("{", "[")):
                try:
                    json.loads(snip)
                    content_type = ContentType.JSON_DATA
                    reasons.append("sniff:json")
                except Exception:
                    pass
            else:
                content_type = ContentType.TEXT
                reasons.append("sniff:text-fallback")
                confidence = 0.6

        # ── 4. Source type from URL ───────────────────────────────────
        if any(p in url for p in _CANVAS_FILE_PATTERNS):
            source_type = SourceType.FILE_DOWNLOAD
            reasons.append("src:canvas-file-download")
        elif host in _CLOUD_DOC_HOSTS:
            source_type = SourceType.CLOUD_DOCUMENT
            reasons.append(f"src:cloud-doc:{host}")
        elif any(h in host for h in _MEDIA_PLATFORMS):
            source_type = SourceType.EMBEDDED_MEDIA
            reasons.append(f"src:media-platform:{host}")
        elif content_type == ContentType.IMAGE:
            source_type = SourceType.IMAGE
        elif content_type in (ContentType.XLSX, ContentType.CSV):
            source_type = SourceType.STRUCTURED_DATASET
        elif content_type == ContentType.HTML:
            source_type = SourceType.WEBSITE_PAGE
        elif content_type in (ContentType.PDF, ContentType.DOCX, ContentType.PPTX, ContentType.TEXT):
            source_type = SourceType.FILE_DOWNLOAD
        else:
            source_type = SourceType.UNKNOWN

        # ── 5. Intent inference ───────────────────────────────────────
        intent = self._infer_intent(content_type, url, content_snippet, context)
        reasons.append(f"intent:{intent.value}")

        if not reasons:
            confidence = 0.3

        logger.debug(f"[classifier] {url[:80]} → ct={content_type.value} src={source_type.value} intent={intent.value} conf={confidence}")
        return Classification(
            source_type=source_type,
            content_type=content_type,
            intent=intent,
            mime_type=normalized_mime,
            confidence=confidence,
            reasons=reasons,
        )

    def _infer_intent(
        self,
        ct: ContentType,
        url: str,
        snippet: str,
        context: Optional[ItemContext],
    ) -> Intent:
        # Check for login walls
        if snippet and re.search(
            r'<input[^>]+type=["\']?password', snippet, re.I
        ):
            return Intent.LOGIN_BARRIER

        if ct in (ContentType.XLSX, ContentType.CSV):
            return Intent.TABULAR

        if ct in (ContentType.PDF, ContentType.DOCX):
            return Intent.NARRATIVE_DOCUMENT

        if ct == ContentType.PPTX:
            return Intent.REFERENCE_MATERIAL

        if ct in (ContentType.AUDIO, ContentType.VIDEO, ContentType.IMAGE):
            return Intent.MEDIA_ASSET

        if ct == ContentType.JSON_DATA:
            return Intent.STRUCTURED_DATA

        if ct == ContentType.HTML:
            # Module / navigation pages
            if any(p in url for p in ("/modules", "/courses/", "/dashboard")):
                if snippet:
                    # High link density → index/navigation
                    link_count = len(re.findall(r"<a\s", snippet, re.I))
                    word_count = len(re.findall(r"\b\w+\b", re.sub(r"<[^>]+>", " ", snippet)))
                    if word_count > 0 and link_count / max(word_count, 1) > 0.12:
                        return Intent.NAVIGATION_PAGE
            if context and context.source_label in ("module", "syllabus"):
                return Intent.REFERENCE_MATERIAL
            return Intent.REFERENCE_MATERIAL

        return Intent.UNKNOWN


# ------------------------------------------------------------------ #
#  Layer 3: Extractors                                                 #
# ------------------------------------------------------------------ #


class BaseExtractor(ABC):
    """All extractors implement this interface."""

    @classmethod
    @abstractmethod
    def can_handle(cls, ct: ContentType) -> bool: ...

    @abstractmethod
    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent: ...

    def _word_count(self, text: str) -> int:
        return len(re.findall(r"\b\w+\b", text))


class HTMLExtractor(BaseExtractor):
    """
    Extracts readable content from HTML pages.
    Strips navigation, ads, and boilerplate. Preserves document structure.
    """

    # CSS selectors for boilerplate to remove
    _NOISE_SELECTORS = [
        "nav", "header", "footer", "aside", ".sidebar", ".menu",
        ".navigation", ".nav", ".breadcrumb", ".breadcrumbs",
        ".ad", ".ads", ".advertisement", ".cookie-banner",
        ".skip-link", "#skip-nav", ".sr-only",
        "script", "style", "noscript",
    ]

    @classmethod
    def can_handle(cls, ct: ContentType) -> bool:
        return ct == ContentType.HTML

    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            return ExtractedContent(
                raw_text=raw if isinstance(raw, str) else raw.decode("utf-8", "ignore"),
                cleaned_text="",
                errors=["beautifulsoup4 not installed"],
            )

        html = raw if isinstance(raw, str) else raw.decode("utf-8", "ignore")
        soup = BeautifulSoup(html, "html.parser")

        # Capture page title
        title_tag = soup.find("title")
        doc_title = title_tag.get_text(strip=True) if title_tag else ""

        # Remove boilerplate
        for sel in self._NOISE_SELECTORS:
            for el in soup.select(sel):
                el.decompose()

        # Extract structured elements
        headings = []
        for tag in soup.find_all(re.compile(r"^h[1-6]$")):
            txt = tag.get_text(separator=" ", strip=True)
            if txt:
                level = int(tag.name[1])
                headings.append({"level": level, "text": txt})

        tables = []
        for tbl in soup.find_all("table"):
            rows = []
            for tr in tbl.find_all("tr"):
                cells = [td.get_text(separator=" ", strip=True)
                         for td in tr.find_all(["td", "th"])]
                if any(cells):
                    rows.append(cells)
            if rows:
                tables.append(rows)

        lists = []
        for ul in soup.find_all(["ul", "ol"]):
            items = [li.get_text(separator=" ", strip=True) for li in ul.find_all("li")]
            items = [i for i in items if i]
            if items:
                lists.extend(items)

        code_blocks = []
        for code in soup.find_all(["code", "pre"]):
            txt = code.get_text()
            if txt.strip():
                code_blocks.append(txt.strip())

        links = []
        for a in soup.find_all("a", href=True):
            href = a.get("href", "")
            text = a.get_text(strip=True)
            if href and text:
                links.append({"href": href, "text": text})

        # Main content extraction — prefer <main>, <article>, <div#content>, fallback to body
        main = (
            soup.find("main")
            or soup.find("article")
            or soup.find(id=re.compile(r"content|main|body", re.I))
            or soup.find(class_=re.compile(r"content|main|body|entry", re.I))
            or soup.body
        )
        raw_text = html
        cleaned_text = (main or soup).get_text(separator="\n", strip=True) if main else ""

        # Detect login wall
        access_status = context.access_status
        if re.search(r'<input[^>]+type=["\']?password', html, re.I):
            access_status = "login_required"
            context.access_status = access_status

        wc = self._word_count(cleaned_text)
        structured = {
            "title": doc_title,
            "headings": headings,
            "tables": tables,
            "lists": lists,
            "code_blocks": code_blocks,
            "links": links,
        }

        return ExtractedContent(
            raw_text=raw_text[:50_000],
            cleaned_text=cleaned_text[:30_000],
            structured=structured,
            word_count=wc,
        )


class PDFExtractor(BaseExtractor):
    """Extracts text and structure from PDF files."""

    @classmethod
    def can_handle(cls, ct: ContentType) -> bool:
        return ct == ContentType.PDF

    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent:
        if isinstance(raw, str):
            raw = raw.encode("utf-8")
        try:
            import pypdf
        except ImportError:
            return ExtractedContent(
                raw_text="", cleaned_text="",
                errors=["pypdf not installed"],
            )

        errors = []
        warnings = []
        pages_text = []
        headings = []
        doc_title = ""
        try:
            reader = pypdf.PdfReader(io.BytesIO(raw))
            doc_title = (reader.metadata or {}).get("/Title", "") or ""
            page_count = len(reader.pages)

            for page_num, page in enumerate(reader.pages):
                try:
                    text = page.extract_text() or ""
                    pages_text.append({"page": page_num + 1, "text": text})
                    # Heuristic: lines in ALL CAPS that are short → headings
                    for line in text.split("\n"):
                        stripped = line.strip()
                        if stripped and stripped == stripped.upper() and 3 < len(stripped) < 80:
                            headings.append({"level": 2, "text": stripped, "page": page_num + 1})
                except Exception as e:
                    warnings.append(f"Page {page_num + 1} extraction error: {e}")

        except Exception as e:
            errors.append(f"PDF parse error: {e}")
            return ExtractedContent(raw_text="", cleaned_text="", errors=errors)

        full_text = "\n\n".join(p["text"] for p in pages_text)
        avg_chars = len(full_text) / max(page_count, 1)
        is_image_heavy = avg_chars < 100  # likely scanned

        if is_image_heavy:
            warnings.append("PDF appears to be image-heavy (scanned). OCR may be needed.")

        structured = {
            "title": doc_title,
            "headings": headings,
            "pages": pages_text,
        }
        return ExtractedContent(
            raw_text=full_text[:50_000],
            cleaned_text=full_text[:30_000],
            structured=structured,
            page_count=page_count,
            word_count=self._word_count(full_text),
            is_image_heavy=is_image_heavy,
            errors=errors,
            warnings=warnings,
        )


class DocxExtractor(BaseExtractor):
    """Extracts text and structure from Word documents."""

    @classmethod
    def can_handle(cls, ct: ContentType) -> bool:
        return ct == ContentType.DOCX

    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent:
        if isinstance(raw, str):
            raw = raw.encode("utf-8")
        try:
            import docx
        except ImportError:
            return ExtractedContent(
                raw_text="", cleaned_text="",
                errors=["python-docx not installed"],
            )
        try:
            doc = docx.Document(io.BytesIO(raw))
        except Exception as e:
            return ExtractedContent(raw_text="", cleaned_text="", errors=[str(e)])

        headings = []
        paragraphs = []
        tables = []

        for para in doc.paragraphs:
            txt = para.text.strip()
            if not txt:
                continue
            style = para.style.name if para.style else ""
            if style.startswith("Heading"):
                try:
                    level = int(style.split()[-1])
                except ValueError:
                    level = 2
                headings.append({"level": level, "text": txt})
            paragraphs.append(txt)

        for tbl in doc.tables:
            rows = []
            for row in tbl.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(cells)
            if rows:
                tables.append(rows)

        full_text = "\n\n".join(paragraphs)
        doc_title = ""
        try:
            core = doc.core_properties
            doc_title = core.title or ""
        except Exception:
            pass

        structured = {
            "title": doc_title,
            "headings": headings,
            "tables": tables,
        }
        return ExtractedContent(
            raw_text=full_text[:50_000],
            cleaned_text=full_text[:30_000],
            structured=structured,
            word_count=self._word_count(full_text),
        )


class SpreadsheetExtractor(BaseExtractor):
    """Extracts schema and sample data from Excel / CSV files."""

    @classmethod
    def can_handle(cls, ct: ContentType) -> bool:
        return ct in (ContentType.XLSX, ContentType.CSV)

    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent:
        # CSV path
        if context.source_url.lower().endswith(".csv") or (
            isinstance(raw, str) and "\n" in raw[:500] and "," in raw[:500]
        ):
            return self._extract_csv(raw)
        return self._extract_xlsx(raw)

    def _extract_csv(self, raw: bytes | str) -> ExtractedContent:
        text = raw if isinstance(raw, str) else raw.decode("utf-8", "ignore")
        try:
            reader = csv.reader(io.StringIO(text))
            rows = list(reader)
        except Exception as e:
            return ExtractedContent(raw_text=text, cleaned_text="", errors=[str(e)])

        headers = rows[0] if rows else []
        sample = rows[1:6] if len(rows) > 1 else []
        schema = {"columns": headers, "row_count": len(rows) - 1, "sample": sample}
        summary = f"CSV: {len(headers)} columns, {len(rows)-1} rows.\nColumns: {', '.join(headers)}"
        structured = {"schema": schema}
        return ExtractedContent(
            raw_text=text[:50_000],
            cleaned_text=summary,
            structured=structured,
            word_count=len(headers),
            sheet_names=["Sheet1"],
        )

    def _extract_xlsx(self, raw: bytes | str) -> ExtractedContent:
        if isinstance(raw, str):
            raw = raw.encode("utf-8")
        try:
            import openpyxl
        except ImportError:
            return ExtractedContent(
                raw_text="", cleaned_text="",
                errors=["openpyxl not installed"],
            )
        try:
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        except Exception as e:
            return ExtractedContent(raw_text="", cleaned_text="", errors=[str(e)])

        sheet_names = wb.sheetnames
        all_text = []
        schemas = {}

        for name in sheet_names:
            ws = wb[name]
            rows = list(ws.iter_rows(values_only=True, max_row=50))
            if not rows:
                continue
            headers = [str(c) if c is not None else "" for c in (rows[0] or [])]
            sample = [[str(c) if c is not None else "" for c in row] for row in rows[1:6]]
            row_count = ws.max_row or 0
            schemas[name] = {"columns": headers, "row_count": row_count, "sample": sample}
            all_text.append(f"Sheet: {name}\nColumns: {', '.join(headers)}")

        summary = "\n\n".join(all_text)
        structured = {"schema": schemas}
        return ExtractedContent(
            raw_text=summary,
            cleaned_text=summary,
            structured=structured,
            sheet_names=sheet_names,
            word_count=sum(len(s.get("columns", [])) for s in schemas.values()),
        )


class PptxExtractor(BaseExtractor):
    """Extracts slide content from PowerPoint presentations."""

    @classmethod
    def can_handle(cls, ct: ContentType) -> bool:
        return ct == ContentType.PPTX

    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent:
        if isinstance(raw, str):
            raw = raw.encode("utf-8")
        try:
            from pptx import Presentation
        except ImportError:
            return ExtractedContent(
                raw_text="", cleaned_text="",
                errors=["python-pptx not installed"],
            )
        try:
            prs = Presentation(io.BytesIO(raw))
        except Exception as e:
            return ExtractedContent(raw_text="", cleaned_text="", errors=[str(e)])

        slides = []
        all_text_parts = []

        for i, slide in enumerate(prs.slides, 1):
            title = ""
            body_parts = []
            notes_text = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if not txt:
                        continue
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PLACEHOLDER
                        pass
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                        ph_idx = shape.placeholder_format.idx
                        if ph_idx == 0:   # title placeholder
                            title = txt
                        else:
                            body_parts.append(txt)
                    else:
                        body_parts.append(txt)

            if slide.has_notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                except Exception:
                    pass

            slide_record = {
                "slide": i,
                "title": title,
                "body": "\n".join(body_parts),
                "notes": notes_text,
            }
            slides.append(slide_record)
            combined = f"Slide {i}: {title}\n{'\n'.join(body_parts)}"
            all_text_parts.append(combined)

        full_text = "\n\n".join(all_text_parts)
        structured = {"slides": slides}
        return ExtractedContent(
            raw_text=full_text[:50_000],
            cleaned_text=full_text[:30_000],
            structured=structured,
            slide_count=len(slides),
            word_count=self._word_count(full_text),
        )


class ImageExtractor(BaseExtractor):
    """
    Stores image metadata and context. No OCR by default.
    Marks as media_asset / archival unless anchor text suggests text content.
    """

    @classmethod
    def can_handle(cls, ct: ContentType) -> bool:
        return ct == ContentType.IMAGE

    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent:
        size = len(raw) if isinstance(raw, bytes) else len(raw.encode())
        anchor = context.anchor_text.lower()
        has_text_content = any(kw in anchor for kw in
                               ("diagram", "chart", "figure", "table", "formula", "equation",
                                "screenshot", "slide", "infographic"))
        warnings = []
        if has_text_content:
            warnings.append("Image may contain text content based on anchor text — consider OCR")

        summary = (
            f"Image from {context.source_url}\n"
            f"Size: {size:,} bytes\n"
            f"Linked from: {context.referring_page}\n"
            f"Anchor text: {context.anchor_text}"
        )
        return ExtractedContent(
            raw_text=summary,
            cleaned_text=summary,
            structured={},
            is_image_heavy=True,
            word_count=0,
            warnings=warnings,
        )


class MediaExtractor(BaseExtractor):
    """
    Handles audio and video. Stores metadata; marks for transcription if no transcript."""

    @classmethod
    def can_handle(cls, ct: ContentType) -> bool:
        return ct in (ContentType.AUDIO, ContentType.VIDEO)

    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent:
        size = len(raw) if isinstance(raw, bytes) else len(raw.encode())
        summary = (
            f"Media ({context.source_url.split('?')[0]})\n"
            f"Size: {size:,} bytes\n"
            f"Linked from: {context.referring_page}"
        )
        context.access_status = "media_requires_transcription"
        return ExtractedContent(
            raw_text=summary,
            cleaned_text=summary,
            structured={},
            word_count=0,
            warnings=["Media file — transcript not yet available"],
        )


class ArchiveExtractor(BaseExtractor):
    """
    Inspects archive contents and records the file list.
    Does not recurse into archives here — records each contained
    file name so they can be queued for individual processing.
    """

    @classmethod
    def can_handle(cls, ct: ContentType) -> bool:
        # Archives arrive as UNKNOWN — we detect by URL extension
        return ct == ContentType.UNKNOWN

    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent:
        url_lower = context.source_url.lower()
        if not any(url_lower.endswith(e) for e in (".zip", ".tar", ".gz", ".7z", ".rar")):
            return RawTextExtractor().extract(raw, context)

        if isinstance(raw, str):
            raw = raw.encode("utf-8", "ignore")

        file_list = []
        errors = []
        try:
            import zipfile
            if zipfile.is_zipfile(io.BytesIO(raw)):
                with zipfile.ZipFile(io.BytesIO(raw)) as zf:
                    file_list = zf.namelist()
        except Exception as e:
            errors.append(f"Archive inspection error: {e}")

        summary = f"Archive: {context.source_url}\nContains {len(file_list)} files:\n" + "\n".join(f"  {f}" for f in file_list[:50])
        structured = {"archive_contents": file_list}
        return ExtractedContent(
            raw_text=summary,
            cleaned_text=summary,
            structured=structured,
            word_count=len(file_list),
            warnings=["Archive — individual files should be re-queued for processing"],
            errors=errors,
        )


class RawTextExtractor(BaseExtractor):
    """
    Safe fallback — stores whatever is available as plain text.
    Always produces a record; never discards.
    """

    @classmethod
    def can_handle(cls, ct: ContentType) -> bool:
        return ct in (ContentType.TEXT, ContentType.JSON_DATA, ContentType.UNKNOWN)

    def extract(self, raw: bytes | str, context: ItemContext) -> ExtractedContent:
        text = raw if isinstance(raw, str) else raw.decode("utf-8", "ignore")
        return ExtractedContent(
            raw_text=text[:50_000],
            cleaned_text=text[:30_000],
            structured={},
            word_count=self._word_count(text),
            warnings=["Processed via RawTextExtractor fallback"],
        )


# ------------------------------------------------------------------ #
#  Layer 2: ItemRouter                                                 #
# ------------------------------------------------------------------ #


class ItemRouter:
    """Maps a Classification to the correct Extractor."""

    # Ordered — first match wins; RawTextExtractor is the universal fallback
    _REGISTRY: list[type[BaseExtractor]] = [
        HTMLExtractor,
        PDFExtractor,
        DocxExtractor,
        SpreadsheetExtractor,
        PptxExtractor,
        ImageExtractor,
        MediaExtractor,
        ArchiveExtractor,
        RawTextExtractor,
    ]

    def route(self, classification: Classification) -> BaseExtractor:
        for cls in self._REGISTRY:
            if cls.can_handle(classification.content_type):
                return cls()
        return RawTextExtractor()


# ------------------------------------------------------------------ #
#  Layer 4: Normalizer                                                 #
# ------------------------------------------------------------------ #

# Tracking/session params to strip from URLs
_STRIP_PARAMS = {
    "utm_source", "utm_medium", "utm_campaign", "utm_term", "utm_content",
    "ref", "referrer", "session", "token", "sid", "fbclid", "gclid",
    "_ga", "mc_cid", "mc_eid",
}


def _canonicalize_url(url: str) -> str:
    """Strip tracking params, normalize scheme/host, remove trailing slashes."""
    try:
        p = urlparse(url)
        qs = parse_qs(p.query, keep_blank_values=False)
        clean_qs = {k: v for k, v in qs.items() if k.lower() not in _STRIP_PARAMS}
        canonical = urlunparse((
            p.scheme.lower(),
            (p.hostname or "").lower() + (f":{p.port}" if p.port else ""),
            p.path.rstrip("/") or "/",
            p.params,
            urlencode(clean_qs, doseq=True),
            "",   # drop fragment
        ))
        return canonical
    except Exception:
        return url


def _recommend_use(classification: Classification, extracted: ExtractedContent) -> str:
    if classification.intent == Intent.LOGIN_BARRIER:
        return "ignore"
    if extracted.word_count < 5 and not extracted.structured_elements if hasattr(extracted, "structured_elements") else not extracted.structured:
        return "archival"
    if classification.intent in (Intent.MEDIA_ASSET,):
        return "archival"
    if classification.intent in (Intent.TABULAR, Intent.STRUCTURED_DATA):
        return "analytics"
    if classification.intent == Intent.NAVIGATION_PAGE:
        return "archival"
    return "search"


class Normalizer:
    """Assembles ExtractedContent + Classification + ItemContext → IntakeRecord."""

    def normalize(
        self,
        extracted: ExtractedContent,
        classification: Classification,
        context: ItemContext,
        extractor: BaseExtractor,
    ) -> IntakeRecord:
        uid = self._make_uid(context.source_url, context.course_id)
        canonical = _canonicalize_url(context.source_url)
        doc_title = (
            extracted.structured.get("title", "")
            or context.document_title
            or context.page_title
            or context.anchor_text
        )
        metadata = {
            "source_url":    context.source_url,
            "canonical_url": canonical,
            "referring_page": context.referring_page,
            "anchor_text":   context.anchor_text,
            "page_title":    context.page_title,
            "title":         doc_title,
            "course_name":   context.course_name,
            "course_id":     context.course_id,
            "module_name":   context.module_name,
            "parent_section": context.parent_section,
            "source_label":  context.source_label,
            "mime_type":     classification.mime_type,
            "language":      context.language,
            "author":        context.author,
            "publisher":     context.publisher,
            "timestamp":     context.timestamp,
            "access_status": context.access_status,
            "word_count":    extracted.word_count,
            "page_count":    extracted.page_count,
            "slide_count":   extracted.slide_count,
            "sheet_names":   extracted.sheet_names,
            "is_image_heavy": extracted.is_image_heavy,
        }
        return IntakeRecord(
            unique_id=uid,
            canonical_source=canonical,
            input_type=classification.source_type.value,
            content_type=classification.content_type.value,
            processing_route=type(extractor).__name__,
            intent=classification.intent.value,
            raw_text=extracted.raw_text,
            cleaned_text=extracted.cleaned_text,
            structured_elements=extracted.structured,
            metadata=metadata,
            quality_score=0.0,   # filled by QualityScorer
            trust_score=0.0,
            dedup_result="new",
            dedup_canonical_id="",
            errors=list(extracted.errors),
            warnings=list(extracted.warnings),
            recommended_use=_recommend_use(classification, extracted),
            decisions=[],
        )

    def _make_uid(self, url: str, course_id: str) -> str:
        h = hashlib.sha256(f"{course_id}:{_canonicalize_url(url)}".encode()).hexdigest()[:20]
        return f"doc_{h}"


# ------------------------------------------------------------------ #
#  Layer 5: Deduplicator                                               #
# ------------------------------------------------------------------ #


class Deduplicator:
    """
    Two-level deduplication:
      1. In-memory set of seen UIDs within the current crawl session
      2. ChromaDB lookup for records persisted from previous crawls
    """

    def __init__(self, kb):
        self.kb = kb
        self._seen_ids: set[str] = set()

    def check(self, record: IntakeRecord) -> tuple[str, str]:
        """Returns (dedup_result, canonical_id_of_original)."""
        # In-memory: same URL seen twice in this crawl
        if record.unique_id in self._seen_ids:
            logger.debug(f"[dedup] in-memory exact dup: {record.unique_id}")
            return "exact_duplicate", record.unique_id

        # Persistent: check if this canonical URL is already in the DB
        existing = None
        try:
            existing = self.kb.find_by_canonical_url(record.canonical_source)
        except Exception as e:
            logger.warning(f"[dedup] KB lookup failed: {e}")

        if existing:
            existing_text = existing.get("cleaned_text", "")
            if existing_text[:200] == record.cleaned_text[:200]:
                logger.debug(f"[dedup] persistent exact dup: {record.canonical_source}")
                return "exact_duplicate", existing["id"]
            logger.debug(f"[dedup] near dup detected: {record.canonical_source}")
            return "near_duplicate", existing["id"]

        self._seen_ids.add(record.unique_id)
        return "new", ""


# ------------------------------------------------------------------ #
#  Layer 6: QualityScorer                                              #
# ------------------------------------------------------------------ #

# Domains whose content should be highly trusted
_HIGH_TRUST_HOSTS = {"instructure.com", "canvaslms.com"}
# Source labels linked directly from Canvas course content → elevated trust
_TRUSTED_SOURCE_LABELS = {"module", "syllabus", "announcements", "assignment"}


class QualityScorer:
    """Assigns quality_score and trust_score (both 0.0–1.0) to an IntakeRecord."""

    def score(self, record: IntakeRecord) -> tuple[float, float]:
        return self._quality(record), self._trust(record)

    def _quality(self, record: IntakeRecord) -> float:
        score = 0.5

        wc = record.metadata.get("word_count", 0)
        if wc > 500:
            score += 0.20
        elif wc > 100:
            score += 0.10
        elif wc < 20:
            score -= 0.25

        if record.structured_elements.get("headings"):
            score += 0.08
        if record.structured_elements.get("tables"):
            score += 0.05
        if record.metadata.get("title"):
            score += 0.05
        if record.metadata.get("is_image_heavy"):
            score -= 0.20
        if record.errors:
            score -= 0.08 * min(len(record.errors), 3)
        if record.dedup_result == "exact_duplicate":
            score -= 0.40
        elif record.dedup_result == "near_duplicate":
            score -= 0.15
        if record.intent == Intent.NAVIGATION_PAGE.value:
            score -= 0.10
        if record.intent == Intent.LOGIN_BARRIER.value:
            score = 0.05  # near-zero — blocked content

        return round(max(0.0, min(1.0, score)), 3)

    def _trust(self, record: IntakeRecord) -> float:
        try:
            host = urlparse(record.canonical_source).hostname or ""
        except Exception:
            host = ""

        if any(h in host for h in _HIGH_TRUST_HOSTS):
            return 0.95
        if record.metadata.get("source_label") in _TRUSTED_SOURCE_LABELS:
            return 0.85
        if record.metadata.get("course_id"):
            return 0.75   # came from a course crawl even if external
        return 0.60


# ------------------------------------------------------------------ #
#  Layer 7: DecisionLogger                                             #
# ------------------------------------------------------------------ #


class DecisionLogger:
    """Appends a structured decision entry to the record and emits a log line."""

    def log_step(
        self,
        record: IntakeRecord,
        step: str,
        decision: str,
        reason: str,
        fallback_used: bool = False,
    ) -> None:
        entry = {
            "step": step,
            "decision": decision,
            "reason": reason,
            "fallback_used": fallback_used,
            "timestamp": datetime.now(timezone.utc).isoformat(),
        }
        record.decisions.append(entry)
        level = logging.WARNING if fallback_used else logging.DEBUG
        logger.log(level, f"[intake:{step}] {decision} | {reason}"
                   + (" [FALLBACK]" if fallback_used else ""))


# ------------------------------------------------------------------ #
#  Pipeline Orchestrator                                               #
# ------------------------------------------------------------------ #


class IntakePipeline:
    """
    Orchestrates all seven layers for a single item.

    Usage:
        pipeline = IntakePipeline(kb)
        record = await pipeline.process(html_bytes, context, "text/html")
        kb.store_intake_record(record)
    """

    def __init__(self, kb):
        self.kb = kb
        self.classifier = ItemClassifier()
        self.router = ItemRouter()
        self.normalizer = Normalizer()
        self.deduplicator = Deduplicator(kb)
        self.scorer = QualityScorer()
        self.decision_logger = DecisionLogger()

    async def process(
        self,
        content: bytes | str,
        context: ItemContext,
        mime_type_hint: str = "",
    ) -> IntakeRecord:
        """
        Run the full pipeline for one item.
        Never raises — on any unhandled error falls through to a
        minimal fallback record so nothing is silently discarded.
        """
        snippet = ""
        try:
            snippet = (
                content[:512]
                if isinstance(content, str)
                else content[:512].decode("utf-8", "ignore")
            )
        except Exception:
            pass

        try:
            # ── Step 1: Classify ─────────────────────────────────────
            classification = self.classifier.classify(
                context.source_url, mime_type_hint, snippet, context
            )

            # ── Step 2: Route ────────────────────────────────────────
            extractor = self.router.route(classification)

            # ── Step 3: Extract ──────────────────────────────────────
            extracted = extractor.extract(content, context)

            # ── Step 4: Normalize ────────────────────────────────────
            record = self.normalizer.normalize(extracted, classification, context, extractor)

            # ── Step 5: Deduplicate ──────────────────────────────────
            dedup_result, canonical_id = self.deduplicator.check(record)
            record.dedup_result = dedup_result
            record.dedup_canonical_id = canonical_id

            # ── Step 6: Score ────────────────────────────────────────
            record.quality_score, record.trust_score = self.scorer.score(record)

            # ── Step 7: Log decisions ────────────────────────────────
            self.decision_logger.log_step(
                record, "classify", classification.content_type.value,
                " | ".join(classification.reasons),
            )
            self.decision_logger.log_step(
                record, "route", type(extractor).__name__,
                f"content_type={classification.content_type.value}",
            )
            self.decision_logger.log_step(
                record, "dedup", dedup_result,
                f"canonical_id={canonical_id}" if canonical_id else "unique",
            )
            self.decision_logger.log_step(
                record, "score",
                f"quality={record.quality_score} trust={record.trust_score}",
                f"word_count={record.metadata.get('word_count',0)} "
                f"intent={record.intent}",
            )
            return record

        except Exception as exc:
            logger.error(
                f"[intake] Unhandled pipeline error for {context.source_url}: {exc}",
                exc_info=True,
            )
            return self._make_fallback_record(context, str(exc))

    def _make_fallback_record(self, context: ItemContext, error: str) -> IntakeRecord:
        """
        Produces a minimal record on catastrophic failure.
        Never raises — this is the last safety net.
        """
        uid = "doc_" + hashlib.sha256(context.source_url.encode()).hexdigest()[:20]
        record = IntakeRecord(
            unique_id=uid,
            canonical_source=_canonicalize_url(context.source_url),
            input_type=SourceType.UNKNOWN.value,
            content_type=ContentType.UNKNOWN.value,
            processing_route="fallback",
            intent=Intent.UNKNOWN.value,
            raw_text="",
            cleaned_text="",
            structured_elements={},
            metadata={
                "source_url":   context.source_url,
                "course_name":  context.course_name,
                "course_id":    context.course_id,
                "source_label": context.source_label,
                "access_status": "unknown",
                "word_count":   0,
            },
            quality_score=0.05,
            trust_score=0.30,
            dedup_result="new",
            dedup_canonical_id="",
            errors=[error],
            warnings=["Processed via emergency fallback — human review recommended"],
            recommended_use="ignore",
            decisions=[{
                "step": "pipeline",
                "decision": "fallback",
                "reason": error,
                "fallback_used": True,
                "timestamp": datetime.now(timezone.utc).isoformat(),
            }],
        )
        self.decision_logger.log_step(
            record, "pipeline", "fallback", error, fallback_used=True
        )
        return record
