"""
Document Ingester
Downloads and extracts readable text from every file and link attached to
Canvas pages, modules, and assignments.

Handles three cases:
  Case 1: Canvas-hosted files (PDF, Word, PowerPoint, plain text)
          Downloaded via Playwright's API client (shares authenticated cookies,
          avoids slow CDP integer-array serialisation).
  Case 2: Embedded / linked documents (Google Docs, OneDrive, SharePoint)
          Fetched via httpx using session cookies where possible.
  Case 3: External platform links (VitalSource, Pearson, McGraw-Hill, etc.)
          Cannot be automated — flagged so the user can upload manually.

Review fixes applied:
  - RF-Ingester-bypass : accepts goto_fn callable so rate limiting is shared
                         with the crawler, not bypassed
  - RF-Dedup           : global _seen_urls set prevents the same URL from being
                         processed by multiple discovery passes
  - RF-Recursion       : global dedup also prevents page-embed infinite loops
  - RF-CDP-bytes       : uses page.context.request.get() instead of CDP
                         integer-array transfer for file downloads
  - RF-httpx-import    : httpx imported at module level, not inside methods
  - RF-FileSizeLimit   : HEAD request checks Content-Length before downloading
"""

import asyncio
import io
import logging
import os
import re
import tempfile
from pathlib import Path
from typing import Callable, Coroutine, Optional
from urllib.parse import urlparse

import httpx

logger = logging.getLogger(__name__)

# ── External platform detection ──────────────────────────────────────

EXTERNAL_PLATFORMS = {
    "vitalsource.com":   "VitalSource",
    "chegg.com":         "Chegg",
    "pearson.com":       "Pearson",
    "mheducation.com":   "McGraw-Hill",
    "cengage.com":       "Cengage",
    "oup.com":           "Oxford University Press",
    "sagepub.com":       "SAGE",
    "wiley.com":         "Wiley",
    "taylorfrancis.com": "Taylor & Francis",
    "springer.com":      "Springer",
    "elsevier.com":      "Elsevier",
    "amazon.com":        "Amazon (textbook)",
}

GOOGLE_DOCS_HOSTS = {"docs.google.com", "drive.google.com"}
MICROSOFT_HOSTS   = {"onedrive.live.com", "sharepoint.com", "1drv.ms"}
YOUTUBE_HOSTS     = {"youtube.com", "youtu.be", "www.youtube.com"}

# Hard cap on individual file downloads to protect against OOM
MAX_FILE_BYTES = 15 * 1024 * 1024  # 15 MB

# Domains whose URLs appear as module items but contain no course-specific content
# (generic Canvas help pages, university support portals, etc.)
_SKIP_MODULE_DOMAINS = {
    "community.instructure.com",
    "community.canvaslms.com",
    "wilmu.edu",  # general university links, not course material
}


def classify_url(url: str) -> str:
    """
    Return one of: 'canvas_file' | 'google_doc' | 'microsoft_doc' |
                   'youtube' | 'external_platform' | 'web_page' | 'unknown'
    """
    if not url:
        return "unknown"
    parsed = urlparse(url.lower())
    host = parsed.netloc.replace("www.", "")

    for domain in EXTERNAL_PLATFORMS:
        if domain in host:
            return "external_platform"

    if any(h in host for h in GOOGLE_DOCS_HOSTS):
        return "google_doc"

    if any(h in host for h in MICROSOFT_HOSTS):
        return "microsoft_doc"

    if any(h in host for h in YOUTUBE_HOSTS):
        return "youtube"

    if "/files/" in parsed.path or "/download" in parsed.path:
        return "canvas_file"

    return "web_page"


def detect_external_platform(url: str) -> Optional[str]:
    """Return the platform name if this URL is a known gated platform, else None."""
    host = urlparse(url.lower()).netloc.replace("www.", "")
    for domain, name in EXTERNAL_PLATFORMS.items():
        if domain in host:
            return name
    return None


# ── Text extractors ──────────────────────────────────────────────────

def extract_pdf_text(data: bytes) -> str:
    """Extract all text from a PDF byte stream. Tries pdfplumber first, falls back to pypdf."""
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
        text = "\n\n".join(pages)
        if text.strip():
            return text.strip()
    except Exception as e:
        logger.debug(f"pdfplumber failed: {e}")

    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages).strip()
    except Exception as e:
        logger.debug(f"pypdf failed: {e}")

    return ""


def extract_docx_text(data: bytes) -> str:
    """Extract text from a Word document byte stream."""
    try:
        import docx2txt
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            tmp.write(data)
            tmp_path = tmp.name
        text = docx2txt.process(tmp_path)
        Path(tmp_path).unlink(missing_ok=True)
        return text.strip() if text else ""
    except Exception as e:
        logger.debug(f"docx2txt failed: {e}")

    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.debug(f"python-docx fallback failed: {e}")

    return ""


def extract_pptx_text(data: bytes) -> str:
    """Extract all text from a PowerPoint byte stream."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"Slide {i}:"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            slides.append("\n".join(parts))
        return "\n\n".join(slides).strip()
    except Exception as e:
        logger.debug(f"pptx extraction failed: {e}")
    return ""


def extract_html_text(html: str) -> str:
    """Strip HTML tags and return clean readable text."""
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        lines = [l.strip() for l in soup.get_text(separator="\n").splitlines() if l.strip()]
        return "\n".join(lines)
    except Exception as e:
        logger.debug(f"HTML extraction failed: {e}")
    return html[:5000]


def extract_text_from_bytes(data: bytes, content_type: str, filename: str = "") -> str:
    """Route to the right extractor based on content type and filename."""
    ct = content_type.lower()
    fn = filename.lower()

    if "pdf" in ct or fn.endswith(".pdf"):
        return extract_pdf_text(data)
    if "wordprocessingml" in ct or fn.endswith(".docx"):
        return extract_docx_text(data)
    if "presentationml" in ct or fn.endswith(".pptx"):
        return extract_pptx_text(data)
    if "text/plain" in ct or fn.endswith(".txt"):
        return data.decode("utf-8", errors="replace")
    if "html" in ct:
        return extract_html_text(data.decode("utf-8", errors="replace"))

    text = extract_pdf_text(data)
    if text:
        return text
    text = extract_docx_text(data)
    if text:
        return text

    logger.warning(f"Could not extract text from {filename} ({content_type})")
    return ""


# ── DocumentIngester ─────────────────────────────────────────────────

# Type alias for the rate-limited goto callable passed in from the crawler
GotoFn = Callable[..., Coroutine]


class DocumentIngester:
    """
    Downloads and extracts text from all documents linked in Canvas.
    Uses the Playwright browser session (already authenticated) for Canvas files
    and httpx for external resources.
    """

    def __init__(self, page, base_url: str, goto_fn: Optional[GotoFn] = None):
        """
        page:     Playwright Page object (authenticated Canvas session)
        base_url: Canvas base URL e.g. https://wilmu.instructure.com
        goto_fn:  Rate-limited navigation coroutine (page, url) -> None.
                  If omitted, falls back to a plain page.goto with a 1s delay.
                  RF-Ingester-bypass: always pass this from CanvasCrawler._polite_goto
                  so the same rate-limiting applies to document ingestion.
        """
        self.page = page
        self.base_url = base_url
        self._goto_fn = goto_fn
        self.results: list[dict] = []
        self.flagged: list[dict] = []
        # RF-Dedup + RF-Recursion: single set shared across all discovery passes
        self._seen_urls: set[str] = set()

    async def _goto(self, url: str) -> None:
        """Navigate using the provided rate-limited helper, or a safe fallback."""
        if self._goto_fn:
            await self._goto_fn(self.page, url)
        else:
            await asyncio.sleep(1.0)
            await self.page.goto(url, wait_until="networkidle")

    # ------------------------------------------------------------------ #
    #  Public entry point                                                  #
    # ------------------------------------------------------------------ #

    async def ingest_course_documents(self, course_id: str, course_name: str) -> dict:
        """
        Discover and ingest all documents for one course.
        Returns dict with 'ingested' (list of text results) and 'flagged' (external links).
        """
        logger.debug(f"[ingest_course_documents] Starting for course: {course_name} (id={course_id})")
        self.results = []
        self.flagged = []
        self._seen_urls = set()  # reset per course

        # Phase 1: Module items — primary source of truth.
        # Files, assignments, discussions, pages, and external links are all
        # discovered here with full module context (module name, anchor text).
        # Running this first ensures every item gets its richest possible metadata.
        logger.debug(f"[ingest_course_documents] Phase 1/2 — ingesting Module items")
        await self._ingest_module_items(course_id, course_name)

        # Phase 2: Canvas Pages index — catch-all for instructor-created pages
        # that are not linked from any module. Pages already visited during
        # Phase 1 are in _seen_urls and will be skipped automatically.
        logger.debug(f"[ingest_course_documents] Phase 2/3 — ingesting Canvas Pages (catch-all)")
        await self._ingest_pages(course_id, course_name)

        # Phase 3: Files section — catch-all for files uploaded to the course
        # file browser but not linked from any module or page. Dedup via
        # _seen_urls means anything already downloaded above is silently skipped.
        logger.debug(f"[ingest_course_documents] Phase 3/3 — ingesting Files section (catch-all)")
        await self._ingest_files_page(course_id, course_name)

        logger.info(
            f"Course {course_name}: ingested {len(self.results)} documents, "
            f"flagged {len(self.flagged)} external links"
        )
        logger.debug(f"[ingest_course_documents] Done for {course_name} — seen_urls total: {len(self._seen_urls)}")
        return {"ingested": self.results, "flagged": self.flagged}

    # ------------------------------------------------------------------ #
    #  Discovery methods                                                   #
    # ------------------------------------------------------------------ #

    async def _ingest_files_page(self, course_id: str, course_name: str) -> None:
        """Download every file listed in the Canvas Files section."""
        logger.debug(f"[_ingest_files_page] Navigating to files for {course_name}")
        try:
            await self._goto(f"{self.base_url}/courses/{course_id}/files")
            try:
                await self.page.wait_for_selector(
                    "a.ef-name-col__link, tr.ef-item-row, .ef-directory", timeout=10000
                )
            except Exception:
                logger.debug(f"[_ingest_files_page] Selector timeout — falling back to 3s wait")
                await self.page.wait_for_timeout(3000)

            file_links = await self.page.query_selector_all(
                "a.ef-name-col__link, tr.ef-item-row a[href*='/files/'], a[href*='/download']"
            )
            logger.debug(f"[_ingest_files_page] Found {len(file_links)} file links for {course_name}")

            seen_hrefs: set[str] = set()
            for link in file_links:
                href = await link.get_attribute("href") or ""
                name = (await link.inner_text()).strip()
                if not href or href in seen_hrefs:
                    continue
                seen_hrefs.add(href)
                full_url = href if href.startswith("http") else f"{self.base_url}{href}"
                logger.debug(f"[_ingest_files_page] Processing file: {name!r} -> {full_url}")
                await self._process_url(full_url, name, course_name, source="files")

        except Exception as e:
            logger.warning(f"Error ingesting files page for {course_name}: {e}")

    async def _ingest_pages(self, course_id: str, course_name: str) -> None:
        """Read all instructor-created Canvas pages for embedded links and content."""
        logger.debug(f"[_ingest_pages] Navigating to pages index for {course_name}")
        try:
            await self._goto(f"{self.base_url}/courses/{course_id}/pages")
            try:
                await self.page.wait_for_selector(
                    "a.wiki-page-link, .pages-index, table.index_content", timeout=8000
                )
            except Exception:
                logger.debug(f"[_ingest_pages] Selector timeout — falling back to 2s wait")
                await self.page.wait_for_timeout(2000)

            page_links = await self.page.query_selector_all(
                "a.wiki-page-link, "
                "table.index_content a[href*='/pages/'], "
                "a[href*='/courses/'][href*='/pages/']"
            )
            logger.debug(f"[_ingest_pages] Found {len(page_links)} Canvas page links for {course_name}")

            seen_hrefs: set[str] = set()
            for link in page_links:
                href = await link.get_attribute("href") or ""
                if not href or href in seen_hrefs:
                    continue
                # Skip template/placeholder hrefs
                if "{{" in href or "}}" in href:
                    continue
                seen_hrefs.add(href)
                full_url = href if href.startswith("http") else f"{self.base_url}{href}"
                logger.debug(f"[_ingest_pages] Scraping Canvas page: {full_url}")
                # RF-Dedup: checked inside _scrape_canvas_page
                await self._scrape_canvas_page(full_url, course_name)

        except Exception as e:
            logger.warning(f"Error ingesting pages for {course_name}: {e}")

    async def _ingest_module_items(self, course_id: str, course_name: str) -> None:
        """
        Walk all module items and process every linked resource.

        Iterates module-by-module so we can attach the module name to each
        item for context. Processes ALL item types — files, Canvas pages,
        assignments, discussions, external URLs — letting _process_url()
        decide the right extraction path.
        """
        logger.debug(f"[_ingest_module_items] Navigating to modules for {course_name}")
        try:
            await self._goto(f"{self.base_url}/courses/{course_id}/modules")
            try:
                await self.page.wait_for_selector(".context_module_item", timeout=8000)
            except Exception:
                logger.debug(f"[_ingest_module_items] Selector timeout — falling back to 2s wait")
                await self.page.wait_for_timeout(2000)
            # Extra settle time for React to finish populating item hrefs.
            # Without this, ExternalTool items have {{ id }} Handlebars placeholders.
            await self.page.wait_for_timeout(3000)

            # Iterate per module so we capture module names for context
            modules = await self.page.query_selector_all(".context_module")
            if not modules:
                # Fallback: no module wrappers found — process all items flat
                modules = [None]

            total_items = 0
            for module_el in modules:
                # Get module name
                module_name = ""
                if module_el:
                    try:
                        name_el = await module_el.query_selector(".ig-header-title")
                        if name_el:
                            module_name = (await name_el.inner_text()).strip()
                    except Exception:
                        pass

                # Get items within this module (or all items on the page for flat fallback)
                item_scope = module_el if module_el else self.page
                items = await item_scope.query_selector_all(".context_module_item")
                logger.debug(
                    f"[_ingest_module_items] Module {module_name!r}: {len(items)} items"
                )

                for item in items:
                    link = await item.query_selector("a.external_url_link, a[href*='external']")
                    if not link:
                        link = await item.query_selector("a.title")
                    if not link:
                        continue

                    href = await link.get_attribute("href") or ""
                    title = (await link.inner_text()).strip()

                    if not href:
                        continue

                    # Skip non-navigable schemes (mailto, javascript, tel, anchors)
                    if href.startswith(("mailto:", "javascript:", "tel:", "#")):
                        continue

                    # Skip unrendered Handlebars template hrefs
                    if "{{" in href or "}}" in href:
                        logger.debug(f"Skipping unrendered module item href: {href!r}")
                        continue

                    full_url = href if href.startswith("http") else f"{self.base_url}{href}"

                    # Skip generic help/support domains — no course-specific content
                    parsed_host = urlparse(full_url).netloc.replace("www.", "")
                    if any(d in parsed_host for d in _SKIP_MODULE_DOMAINS):
                        logger.debug(f"[_ingest_module_items] Skipping noise domain: {full_url}")
                        continue

                    logger.debug(
                        f"[_ingest_module_items] Processing: {title!r} "
                        f"(module={module_name!r}) -> {full_url}"
                    )
                    await self._process_url(
                        full_url, title, course_name,
                        source="module", module_name=module_name
                    )
                    total_items += 1

            logger.info(f"[_ingest_module_items] Processed {total_items} module items for {course_name}")

        except Exception as e:
            logger.warning(f"Error ingesting module items for {course_name}: {e}")

    # ------------------------------------------------------------------ #
    #  Processing                                                          #
    # ------------------------------------------------------------------ #

    async def _process_url(
        self,
        url: str,
        title: str,
        course_name: str,
        source: str,
        module_name: str = "",
    ) -> None:
        """
        Classify a URL and route it to the appropriate handler.

        Routing priority:
          1. External gated platforms (VitalSource, Pearson…) → flag for manual upload
          2. Canvas-hosted files (/files/, /download) → download + extract
          3. Google Docs / Drive → export-URL fetch
          4. Microsoft OneDrive / SharePoint → browser fetch
          5. YouTube → transcript API
          6. Any other Canvas-internal URL (assignments, discussions, pages,
             module items, etc.) → browser scrape
          7. Generic web page → HTTP fetch + HTML extraction

        RF-Dedup + RF-Recursion: global _seen_urls guard prevents re-processing
        the same URL across multiple discovery passes.
        """
        if url in self._seen_urls:
            logger.debug(f"[_process_url] Skipping duplicate URL: {url}")
            return

        url_type = classify_url(url)
        logger.debug(
            f"[_process_url] {url_type.upper()} | {title!r} | "
            f"source={source} module={module_name!r} | {url}"
        )

        if url_type == "external_platform":
            platform = detect_external_platform(url)
            logger.info(f"Flagging external platform: {platform} -- {title}")
            self.flagged.append({
                "title": title,
                "url": url,
                "platform": platform,
                "course_name": course_name,
                "source": source,
                "module_name": module_name,
                "note": f"Manual upload required -- {platform} requires separate login",
            })

        elif url_type == "canvas_file":
            await self._download_canvas_file(url, title, course_name, source)

        elif url_type in ("google_doc", "microsoft_doc"):
            await self._fetch_embedded_doc(url, title, course_name, source)

        elif url_type == "youtube":
            await self._fetch_youtube_transcript(url, title, course_name, source)

        elif self.base_url and url.startswith(self.base_url):
            # Any Canvas-internal URL that isn't a file download:
            # assignments, discussions, module item redirects, wiki pages, etc.
            # Use the browser (already authenticated) to scrape readable content.
            await self._scrape_canvas_page(
                url, course_name, title=title, source=source, module_name=module_name
            )

        else:
            # Generic external web page
            await self._fetch_web_page(url, title, course_name, source)

        # Mark URL as seen for all non-canvas-page paths.
        # Canvas-internal pages are marked inside _scrape_canvas_page itself
        # (so that the guard there works for direct calls from _ingest_pages).
        if url not in self._seen_urls:
            self._seen_urls.add(url)

    async def _scrape_canvas_page(
        self,
        url: str,
        course_name: str,
        title: str = "",
        source: str = "page",
        module_name: str = "",
    ) -> None:
        """
        Scrape any Canvas-internal page: wiki pages, assignments, discussions,
        module item redirect URLs, etc.

        Uses a priority list of content selectors so the right body is captured
        regardless of the page type. After extracting text it also follows any
        embedded external links so resources linked from within a page are also
        ingested.

        RF-Recursion: _seen_urls prevents cycles when pages link each other.
        """
        # _seen_urls guard is in _process_url — but guard again here for calls
        # that come directly from _ingest_pages (which uses _scrape_canvas_page
        # directly without going through _process_url).
        if url in self._seen_urls:
            logger.debug(f"[_scrape_canvas_page] Skipping duplicate: {url}")
            return
        self._seen_urls.add(url)

        logger.debug(f"[_scrape_canvas_page] Navigating to: {url}")
        try:
            await self._goto(url)
            await asyncio.sleep(1.0)  # slightly longer for React-heavy Canvas pages

            # Save snapshot so module item visits are visible in debug_snapshots/
            try:
                snap_dir = Path(os.environ.get("DATA_DIR", "data")) / "debug_snapshots"
                snap_dir.mkdir(parents=True, exist_ok=True)
                safe_label = re.sub(r"[^\w\-]", "_", f"module_item_{title or url.split('/')[-1]}")[:80]
                (snap_dir / f"{safe_label}.html").write_text(
                    await self.page.content(), encoding="utf-8", errors="replace"
                )
            except Exception as _snap_err:
                logger.debug(f"[_scrape_canvas_page] Snapshot save failed: {_snap_err}")

            # Resolve title from page if not provided
            if not title:
                for sel in ["h1.page-title", "h1.title", ".discussion-title h1", "h1"]:
                    title_el = await self.page.query_selector(sel)
                    if title_el:
                        title = (await title_el.inner_text()).strip()
                        break
                if not title:
                    title = "Canvas Page"
            logger.debug(f"[_scrape_canvas_page] Page title: {title!r}")

            # Content selectors in priority order — covers wiki pages, assignments,
            # discussions, announcements, and generic content areas
            content_selectors = [
                ".show-content",                  # wiki/pages
                "#wiki_page_show",                # wiki (older Canvas)
                ".assignment-description",        # assignment body
                ".description.user_content",      # assignment alt
                "#assignment_description",        # assignment alt 2
                ".discussion-entries",            # discussion replies
                ".discussion-section .entry-content",
                ".announcement-content",          # announcements
                "#announcement_message_holder",
                "#content",                       # generic content
                "main",                           # HTML5 main
            ]

            body_el = None
            for sel in content_selectors:
                body_el = await self.page.query_selector(sel)
                if body_el:
                    logger.debug(f"[_scrape_canvas_page] Content selector matched: {sel}")
                    break

            if not body_el:
                logger.debug(f"[_scrape_canvas_page] No specific content selector matched — using body")
                body_el = await self.page.query_selector("body")

            text = ""
            if body_el:
                text = (await body_el.inner_text()).strip()
                logger.debug(f"[_scrape_canvas_page] Body text: {len(text)} chars")

            if text and len(text) > 50:
                # Vision: supplement with screenshot if rich media is present
                imgs = await (body_el or self.page).query_selector_all("img")
                tables = await (body_el or self.page).query_selector_all("table")
                logger.debug(f"[_scrape_canvas_page] {len(imgs)} images, {len(tables)} tables")
                if imgs or tables:
                    try:
                        from agent.brain import describe_page_visuals
                        screenshot = await self.page.screenshot(full_page=True)
                        logger.debug(f"[_scrape_canvas_page] Screenshot {len(screenshot):,} bytes")
                        vision_text = await describe_page_visuals(screenshot, url)
                        if vision_text:
                            logger.debug(f"[_scrape_canvas_page] Vision: {len(vision_text)} chars")
                            text = text + f"\n\n[VISUAL CONTENT]\n{vision_text}"
                    except Exception as e:
                        logger.debug(f"Vision screenshot skipped: {e}")

                # Determine doc_type from URL path
                url_path = urlparse(url).path.lower()
                if "/assignments/" in url_path:
                    doc_type = "canvas_assignment"
                elif "/discussion_topics/" in url_path:
                    doc_type = "canvas_discussion"
                elif "/announcements/" in url_path:
                    doc_type = "canvas_announcement"
                elif "/pages/" in url_path:
                    doc_type = "canvas_page"
                else:
                    doc_type = "canvas_content"

                from agent.brain import enrich_for_knowledge_base
                enriched = await enrich_for_knowledge_base(
                    text, title, course_name, doc_type, url
                )
                self._store_result(
                    title, enriched, url, course_name,
                    source=source, doc_type=doc_type,
                    module_name=module_name,
                )
                logger.info(
                    f"[_scrape_canvas_page] Stored {len(enriched):,} chars: "
                    f"{title!r} ({doc_type}) module={module_name!r}"
                )
            else:
                logger.debug(f"[_scrape_canvas_page] Insufficient text ({len(text)} chars) — skipping storage")

            # Follow embedded links — Canvas-internal and Canvas files only.
            # Per spec §12 (Boilerplate Suppression), external web pages embedded in
            # Canvas pages (help links, university admin sites, etc.) are not course
            # content and must be skipped. Only follow links that resolve to
            # Canvas-internal URLs, Canvas file downloads, or gated external platforms
            # (which get flagged for manual upload, not fetched).
            link_els = await self.page.query_selector_all(
                ".show-content a[href], #wiki_page_show a[href], "
                ".assignment-description a[href], .discussion-entries a[href], "
                ".announcement-content a[href], #content a[href]"
            )
            logger.debug(f"[_scrape_canvas_page] {len(link_els)} embedded links on {title!r}")
            for link_el in link_els:
                href = await link_el.get_attribute("href") or ""
                link_title = (await link_el.inner_text()).strip() or title
                if not href or href.startswith(("#", "mailto:", "javascript:", "tel:")):
                    continue
                full_url = href if href.startswith("http") else f"{self.base_url}{href}"
                # Only follow Canvas-internal links or known external platforms from page content.
                # Generic external web pages are boilerplate (student services, help, etc.).
                if not full_url.startswith(self.base_url):
                    parsed_host = urlparse(full_url).netloc.replace("www.", "")
                    is_external_platform = any(d in parsed_host for d in EXTERNAL_PLATFORMS)
                    if not is_external_platform:
                        logger.debug(f"[_scrape_canvas_page] Skipping external link: {full_url}")
                        continue
                await self._process_url(
                    full_url, link_title, course_name,
                    source="page_embed", module_name=module_name
                )

        except Exception as e:
            logger.warning(f"Error scraping Canvas page {url}: {e}")

    async def _download_canvas_file(
        self, url: str, title: str, course_name: str, source: str
    ) -> None:
        """
        Download a Canvas-hosted file using Playwright's API client.

        RF-CDP-bytes: uses page.context.request.get() instead of JS
        Array.from(new Uint8Array(buffer)) via CDP. The CDP approach
        serialises the entire file as a JSON integer array, which is
        ~3x the file size in memory and extremely slow for PDFs > a few MB.
        Playwright's API client shares the browser's authenticated cookies
        without going through the CDP protocol.

        RF-FileSizeLimit: a HEAD request checks Content-Length before
        downloading, skipping files larger than MAX_FILE_BYTES (15 MB).
        """
        try:
            logger.info(f"Downloading Canvas file: {title}")

            # RF-FileSizeLimit: check size before committing to download.
            # timeout=10000 ms prevents indefinite hangs on slow Canvas responses.
            head_resp = await self.page.context.request.fetch(
                url, method="HEAD", timeout=10000
            )
            content_length = int(head_resp.headers.get("content-length", "0") or "0")
            if content_length > MAX_FILE_BYTES:
                logger.warning(
                    f"Skipping {title}: file is {content_length / 1024 / 1024:.1f} MB "
                    f"(limit {MAX_FILE_BYTES // 1024 // 1024} MB)"
                )
                return

            # RF-CDP-bytes: Playwright API client download — no CDP overhead.
            # RF-DownloadTimeout: 30 s ceiling so a stalled file server can't
            # hang the entire crawl indefinitely.
            response = await self.page.context.request.get(url, timeout=30000)
            if not response.ok:
                logger.warning(f"Failed to download {title}: HTTP {response.status}")
                return

            data = await response.body()
            if not data:
                logger.warning(f"Empty file downloaded: {title}")
                return

            content_type = response.headers.get("content-type", "")
            filename = Path(urlparse(url).path).name
            text = extract_text_from_bytes(data, content_type, filename)

            if text:
                from agent.brain import enrich_for_knowledge_base
                enriched = await enrich_for_knowledge_base(
                    text, title, course_name, "canvas_file", url
                )
                self._store_result(title, enriched, url, course_name, source, "canvas_file")
                logger.info(f"Extracted {len(text):,} chars from: {title}")
            else:
                logger.warning(f"No text extracted from: {title} ({content_type})")

        except Exception as e:
            logger.warning(f"Error downloading Canvas file {title}: {e}", exc_info=True)

    async def _fetch_embedded_doc(
        self, url: str, title: str, course_name: str, source: str
    ) -> None:
        """Attempt to fetch a Google Doc or Microsoft document."""
        try:
            if "docs.google.com" in url:
                doc_id_match = re.search(r"/d/([a-zA-Z0-9_-]+)", url)
                if doc_id_match:
                    doc_id = doc_id_match.group(1)
                    export_url = (
                        f"https://docs.google.com/document/d/{doc_id}/export?format=txt"
                    )
                    text = await self._http_get_text(export_url)
                    if text:
                        from agent.brain import enrich_for_knowledge_base
                        enriched = await enrich_for_knowledge_base(
                            text, title, course_name, "google_doc", url
                        )
                        self._store_result(title, enriched, url, course_name, source, "google_doc")
                        return

            if "drive.google.com" in url:
                file_id_match = re.search(r"/file/d/([a-zA-Z0-9_-]+)", url)
                if file_id_match:
                    file_id = file_id_match.group(1)
                    export_url = (
                        f"https://drive.google.com/uc?export=download&id={file_id}"
                    )
                    data, ct = await self._http_get_bytes(export_url)
                    if data:
                        text = extract_text_from_bytes(data, ct)
                        if text:
                            from agent.brain import enrich_for_knowledge_base
                            enriched = await enrich_for_knowledge_base(
                                text, title, course_name, "google_drive", url
                            )
                            self._store_result(
                                title, enriched, url, course_name, source, "google_drive"
                            )
                            return

            if any(h in url for h in MICROSOFT_HOSTS):
                await self._goto(url)
                await asyncio.sleep(2.0)
                body_el = await self.page.query_selector("body")
                if body_el:
                    text = await body_el.inner_text()
                    if len(text.strip()) > 100:
                        from agent.brain import enrich_for_knowledge_base
                        enriched = await enrich_for_knowledge_base(
                            text.strip(), title, course_name, "microsoft_doc", url
                        )
                        self._store_result(
                            title, enriched, url, course_name, source, "microsoft_doc"
                        )
                        return

            logger.info(f"Could not auto-extract embedded doc: {title} -- flagging")
            self.flagged.append({
                "title": title,
                "url": url,
                "platform": "Embedded Document",
                "course_name": course_name,
                "source": source,
                "note": "Could not auto-extract -- upload manually if critical",
            })

        except Exception as e:
            logger.warning(f"Error fetching embedded doc {title}: {e}")

    async def _fetch_web_page(
        self, url: str, title: str, course_name: str, source: str
    ) -> None:
        """Fetch and extract text from a generic web page."""
        try:
            text = await self._http_get_text(url)
            if text and len(text) > 200:
                from agent.brain import enrich_for_knowledge_base
                enriched = await enrich_for_knowledge_base(
                    text[:50000], title, course_name, "web_page", url
                )
                self._store_result(title, enriched, url, course_name, source, "web_page")
        except Exception as e:
            logger.debug(f"Could not fetch web page {url}: {e}")

    async def _fetch_youtube_transcript(
        self, url: str, title: str, course_name: str, source: str
    ) -> None:
        """
        Fetch the auto-generated or manual transcript for a YouTube video and
        store it in the knowledge base so the AI can reference spoken content.

        Uses youtube-transcript-api which fetches transcripts without a browser.
        Falls back to flagging the video if no transcript is available.
        """
        try:
            from urllib.parse import parse_qs
            from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound, TranscriptsDisabled

            # Extract video ID from both youtube.com/watch?v=ID and youtu.be/ID forms
            parsed = urlparse(url)
            if "youtu.be" in parsed.netloc:
                video_id = parsed.path.lstrip("/").split("?")[0]
            else:
                qs = parse_qs(parsed.query)
                video_id = (qs.get("v") or [""])[0]

            if not video_id:
                logger.warning(f"Could not extract video ID from YouTube URL: {url}")
                return

            logger.info(f"Fetching YouTube transcript: {title} ({video_id})")

            loop = asyncio.get_running_loop()
            transcript_list = await loop.run_in_executor(
                None, lambda: YouTubeTranscriptApi.get_transcript(video_id)
            )

            # Join all transcript segments into readable text
            transcript_text = " ".join(
                segment.get("text", "") for segment in transcript_list
            ).strip()

            if not transcript_text:
                logger.warning(f"Empty transcript for YouTube video: {title}")
                return

            full_text = f"YouTube Video: {title}\nURL: {url}\n\nTranscript:\n{transcript_text}"

            from agent.brain import enrich_for_knowledge_base
            enriched = await enrich_for_knowledge_base(
                full_text, title, course_name, "youtube_video", url
            )
            self._store_result(title, enriched, url, course_name, source, "youtube_video")
            logger.info(f"YouTube transcript stored: {title} ({len(transcript_text):,} chars)")

        except Exception as e:
            logger.info(f"YouTube transcript unavailable for {url}: {e}")
            # Store a reference entry so the AI knows the video exists
            self._store_result(
                title,
                f"YouTube Video: {title}\nURL: {url}\nNote: Transcript not available — watch manually.",
                url,
                course_name,
                source,
                "youtube_video",
            )

    # ------------------------------------------------------------------ #
    #  HTTP helpers                                                        #
    # ------------------------------------------------------------------ #

    async def _http_get_text(self, url: str) -> str:
        """GET a URL and return its text content. RF-httpx-import: httpx at module level."""
        try:
            async with httpx.AsyncClient(timeout=20, follow_redirects=True) as client:
                resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0"})
                if resp.status_code == 200:
                    ct = resp.headers.get("content-type", "")
                    return extract_html_text(resp.text) if "html" in ct else resp.text
        except Exception as e:
            logger.debug(f"HTTP GET failed for {url}: {e}")
        return ""

    async def _http_get_bytes(self, url: str) -> tuple[bytes, str]:
        """GET a URL and return raw bytes plus content-type. RF-httpx-import: module level."""
        try:
            async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
                resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0"})
                if resp.status_code == 200:
                    return resp.content, resp.headers.get("content-type", "")
        except Exception as e:
            logger.debug(f"HTTP bytes GET failed for {url}: {e}")
        return b"", ""

    # ------------------------------------------------------------------ #
    #  Storage                                                             #
    # ------------------------------------------------------------------ #

    def _store_result(
        self,
        title: str,
        text: str,
        url: str,
        course_name: str,
        source: str,
        doc_type: str,
        module_name: str = "",
    ) -> None:
        """
        Append an extracted document to the results list.
        Attaches intake pipeline classification metadata to each result so
        knowledge_base.ingest_knowledge() can store richer ChromaDB metadata.
        """
        # Run synchronous pipeline classification for metadata enrichment
        intake_meta: dict = {}
        try:
            from agent.intake_pipeline import ItemClassifier, ItemContext
            classifier = ItemClassifier()
            ctx = ItemContext(
                source_url=url,
                course_name=course_name,
                module_name=module_name,
                source_label=source,
                anchor_text=title,
            )
            # Classify from URL + first 512 chars of text
            classification = classifier.classify(url, content_snippet=text[:512], context=ctx)
            intake_meta = {
                "content_type": classification.content_type.value,
                "input_type": classification.source_type.value,
                "intent": classification.intent.value,
                "classifier_reasons": " | ".join(classification.reasons),
                "classifier_confidence": classification.confidence,
            }
        except Exception as e:
            logger.debug(f"[_store_result] Classification metadata skipped: {e}")

        self.results.append({
            "title": title,
            "text": text,
            "url": url,
            "course_name": course_name,
            "source": source,
            "doc_type": doc_type,
            "module_name": module_name,
            "char_count": len(text),
            **intake_meta,
        })
