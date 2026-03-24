"""
File Generator
Creates actual submission-ready files from agent-generated content.
Supports: docx, pptx, xlsx, txt

Review fixes applied:
  - RF-Docstrings: docstrings moved to the first statement in every function
                   so Python recognises them (they were after the guard clause)
"""

import os
import re
import logging
from pathlib import Path
from datetime import datetime

logger = logging.getLogger(__name__)

_MD_SEPARATOR_RE = re.compile(r"^\|[\s\-|]+\|$")


def _default_output_dir() -> str:
    """Return the default output directory driven by DATA_DIR env var."""
    data_dir = os.environ.get("DATA_DIR", "data")
    return str(Path(data_dir) / "assignments")


def sanitize_filename(name: str) -> str:
    """Strip characters unsafe for filenames and cap length at 80 chars."""
    return re.sub(r"[^\w\-_. ]", "_", name)[:80]


def generate_docx(content: str, title: str, output_dir: str = "") -> str:
    """
    Generate a Word document from agent-produced markdown-style content.
    Sections marked with ## / ### headers, bullet lists, numbered lists,
    and bold-only lines are all mapped to the correct python-docx styles.
    RF-Docstrings: docstring is now the first statement in the function.
    """
    if not output_dir:
        output_dir = _default_output_dir()

    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    doc = Document()

    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for line in content.split("\n"):
        line = line.strip()

        if not line:
            continue
        if line.startswith("CONFIDENCE:"):
            continue

        if line.startswith("## "):
            doc.add_heading(line[3:], level=1)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("**") and line.endswith("**") and len(line) < 80:
            doc.add_heading(line.strip("*"), level=2)
        elif line.startswith("- ") or line.startswith("* "):
            doc.add_paragraph(line[2:], style="List Bullet")
        elif re.match(r"^\d+\. ", line):
            doc.add_paragraph(re.sub(r"^\d+\. ", "", line), style="List Number")
        else:
            doc.add_paragraph(line)

    section = doc.sections[0]
    section.footer.paragraphs[0].text = f"Generated: {datetime.now().strftime('%Y-%m-%d')}"

    filename = f"{sanitize_filename(title)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    filepath = str(Path(output_dir) / filename)
    doc.save(filepath)
    logger.info(f"Generated DOCX: {filepath}")
    return filepath


def generate_pptx(content: str, title: str, output_dir: str = "") -> str:
    """
    Generate a PowerPoint presentation from agent-produced content.
    Parses 'Slide N: Title / content' format with a markdown-header fallback.
    RF-Docstrings: docstring is now the first statement in the function.
    """
    if not output_dir:
        output_dir = _default_output_dir()

    from pptx import Presentation

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    slides_raw = re.split(r"\nSlide \d+:", "\nSlide 0:" + content)
    slides_data: list[tuple[str, str]] = []

    for raw in slides_raw[1:]:
        lines = raw.strip().split("\n")
        if not lines:
            continue
        slide_title = lines[0].lstrip(":").strip()
        slide_content = "\n".join(lines[1:]).strip()
        slides_data.append((slide_title, slide_content))

    if not slides_data:
        slides_data = [("Title Slide", "")]
        for section in re.split(r"\n(?=#{1,3} |\*\*[A-Z])", content):
            if not section.strip():
                continue
            lines = section.strip().split("\n")
            sec_title = lines[0].lstrip("#").strip("*").strip()
            sec_content = "\n".join(lines[1:]).strip()
            slides_data.append((sec_title, sec_content))

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    subtitle = slide.placeholders.get(1)
    if subtitle:
        subtitle.text = datetime.now().strftime("%B %Y")

    content_layout = prs.slide_layouts[1]
    for slide_title, slide_content in slides_data:
        if not slide_title:
            continue
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_title

        body = slide.placeholders.get(1)
        if body and slide_content:
            tf = body.text_frame
            tf.word_wrap = True

            bullets = [b.strip() for b in slide_content.split("\n") if b.strip()]
            for j, bullet in enumerate(bullets[:8]):
                bullet = re.sub(r"^[-*]\s*", "", bullet)
                if j == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 1 if bullet.startswith("  ") else 0

    filename = f"{sanitize_filename(title)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = str(Path(output_dir) / filename)
    prs.save(filepath)
    logger.info(f"Generated PPTX: {filepath}")
    return filepath


def generate_xlsx(content: str, title: str, output_dir: str = "") -> str:
    """
    Generate an Excel spreadsheet from agent-produced content.
    Parses pipe-separated markdown tables, CSV lines, and plain text rows.
    Markdown separator rows (|---|) are detected and skipped.
    RF-Docstrings: docstring is now the first statement in the function.
    """
    if not output_dir:
        output_dir = _default_output_dir()

    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:31]

    header_font = Font(bold=True, size=12, color="FFFFFF")
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")

    row_num = 1
    for line in content.split("\n"):
        line = line.strip()
        if not line or line.startswith("CONFIDENCE:"):
            continue

        if "|" in line:
            if _MD_SEPARATOR_RE.match(line):
                continue
            cells = [c.strip() for c in line.split("|") if c.strip()]
            for col, cell_val in enumerate(cells, 1):
                cell_obj = ws.cell(row=row_num, column=col, value=cell_val)
                if row_num == 1:
                    cell_obj.font = header_font
                    cell_obj.fill = header_fill
                    cell_obj.alignment = Alignment(horizontal="center")
        elif "," in line:
            for col, cell_val in enumerate(line.split(","), 1):
                ws.cell(row=row_num, column=col, value=cell_val.strip())
        else:
            ws.cell(row=row_num, column=1, value=line)

        row_num += 1

    for col in ws.columns:
        max_len = max(
            (len(str(cell.value)) for cell in col if cell.value),
            default=0,
        )
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 50)

    filename = f"{sanitize_filename(title)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
    filepath = str(Path(output_dir) / filename)
    wb.save(filepath)
    logger.info(f"Generated XLSX: {filepath}")
    return filepath


def generate_text(content: str, title: str, output_dir: str = "") -> str:
    """
    Save plain text content (discussion posts, short answers).
    RF-Docstrings: docstring is now the first statement in the function.
    """
    if not output_dir:
        output_dir = _default_output_dir()

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    filename = f"{sanitize_filename(title)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
    filepath = str(Path(output_dir) / filename)
    Path(filepath).write_text(
        f"{title}\n{'=' * len(title)}\n\n{content}",
        encoding="utf-8",
    )
    logger.info(f"Generated TXT: {filepath}")
    return filepath


def generate_file(content: str, title: str, file_type: str, output_dir: str = "") -> str:
    """
    Route to the correct generator based on file type string.
    RF-Docstrings: docstring is now the first statement in the function.
    """
    if not output_dir:
        output_dir = _default_output_dir()

    ft = file_type.lower().strip(".")
    if ft == "docx":
        return generate_docx(content, title, output_dir)
    if ft == "pptx":
        return generate_pptx(content, title, output_dir)
    if ft in ("xlsx", "xls"):
        return generate_xlsx(content, title, output_dir)
    return generate_text(content, title, output_dir)
